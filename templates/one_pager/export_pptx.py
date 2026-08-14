"""
Export a GT Hub one-pager to an EDITABLE PowerPoint slide (.pptx).

Same YAML source as render.py — this is just the other output format. Use it when
the page needs to be changed by hand or dropped into the Matchmaking deck:

  * every text block is a real text box  -> click and type
  * every image is a real picture        -> right-click > Change Picture
  * the slide is 16:9                    -> paste straight into the deck

Nothing here is a flattened image. Opens in PowerPoint, Keynote, LibreOffice and
Google Slides.

Usage:
    python3 templates/one_pager/export_pptx.py templates/one_pager/data/ligaro.yaml
    python3 templates/one_pager/export_pptx.py templates/one_pager/data/*.yaml --out-dir ~/Desktop
    python3 templates/one_pager/export_pptx.py data/*.yaml --combine deck.pptx   # all pages, one file

Requires: pip3 install python-pptx pillow
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from render import SECTIONS, VISUALS, validate  # noqa: E402  — single source of truth

# ── Layout, in inches. Proportions taken from the ONOX / Arctory reference pages.
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.24
HDR_Y, HDR_H = 0.13, 0.20
RULE_Y = 0.40
CLAIM_Y, CLAIM_H = 0.50, 0.72
BODY_Y = 1.34
BODY_H = SLIDE_H - BODY_Y - 0.26
LEFT_W = 6.40
GAP = 0.17
RIGHT_X = MARGIN + LEFT_W + GAP
RIGHT_W = SLIDE_W - MARGIN - RIGHT_X
VIS_GAP = 0.16
VIS_H = (BODY_H - VIS_GAP) / 2          # 50/50, as in both reference pages
CARD_PAD = 0.17
LOGO_W, LOGO_H = 1.55, 0.50

ACCENT = RGBColor(0x6C, 0x5C, 0xE7)
INK = RGBColor(0x11, 0x11, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAGE_BG = RGBColor(0xF2, 0xF2, 0xF2)
FONT = "Arial"


def _txbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def _run(p, text, *, size, bold=False, italic=False, underline=False, color=INK):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.underline = underline
    r.font.color.rgb = color
    r.font.name = FONT
    return r


def _rect(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    return s


def _fit_cover(src: Path, box_w_in: float, box_h_in: float, tmp_dir: Path) -> Path:
    """Centre-crop `src` to the box aspect ratio so the picture fills its frame
    rather than letterboxing — matching how the reference pages look. Returns a
    path to the cropped copy (the original on disk is never modified)."""
    from PIL import Image

    im = Image.open(src)
    if im.mode in ("RGBA", "LA", "P"):
        im = im.convert("RGBA")
        bg = Image.new("RGB", im.size, (255, 255, 255))
        bg.paste(im, mask=im.split()[-1])
        im = bg
    else:
        im = im.convert("RGB")

    target = box_w_in / box_h_in
    w, h = im.size
    if w / h > target:                       # too wide -> trim the sides
        new_w = int(h * target)
        im = im.crop(((w - new_w) // 2, 0, (w - new_w) // 2 + new_w, h))
    else:                                    # too tall -> trim top/bottom
        new_h = int(w / target)
        im = im.crop((0, (h - new_h) // 2, w, (h - new_h) // 2 + new_h))

    tmp_dir.mkdir(parents=True, exist_ok=True)
    out = tmp_dir / f"{src.stem}_fit.jpg"
    im.save(out, "JPEG", quality=88, optimize=True)
    return out


def _visual(slide, slot: dict, default_label: str, base_dir: Path, tmp_dir: Path, y: float):
    """One image box: a real picture when we have one, otherwise the accent-violet
    placeholder from the blank template, carrying the label so an unfilled slot is
    obviously unfilled."""
    image = slot.get("image")
    label = str(slot.get("label") or default_label)

    if image:
        src = (base_dir / str(image)).resolve()
        if src.exists():
            pic = _fit_cover(src, RIGHT_W, VIS_H, tmp_dir)
            slide.shapes.add_picture(str(pic), Inches(RIGHT_X), Inches(y),
                                     Inches(RIGHT_W), Inches(VIS_H))
            _rect(slide, RIGHT_X, y, RIGHT_W, VIS_H, fill=None, line=INK)
            return
        print(f"      ! image not found, using placeholder: {src}")

    _rect(slide, RIGHT_X, y, RIGHT_W, VIS_H, fill=ACCENT, line=INK)
    tb, tf = _txbox(slide, RIGHT_X + 0.35, y + VIS_H / 2 - 0.45, RIGHT_W - 0.70, 0.90)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, label, size=14, color=WHITE)
    if slot.get("placeholder"):
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
        _run(p2, str(slot["placeholder"]), size=9, color=WHITE)


def build_slide(prs: Presentation, data: dict, base_dir: Path, tmp_dir: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])          # blank layout
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=PAGE_BG)

    meta = data.get("meta") or {}

    # ── header + rule ────────────────────────────────────────────────────────
    tb, tf = _txbox(slide, MARGIN, HDR_Y, 4.0, HDR_H)
    _run(tf.paragraphs[0], "GT Hub", size=8.5)

    tb, tf = _txbox(slide, SLIDE_W - MARGIN - 4.6, HDR_Y, 4.6, HDR_H)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    label = str(meta.get("page_label") or "Matchmaking-Startups")
    num = str(meta.get("page_number") or "")
    _run(p, f"{label}      {num}".rstrip(), size=8.5)

    _rect(slide, MARGIN, RULE_Y, SLIDE_W - 2 * MARGIN, 0.012, fill=INK)

    # ── claim ────────────────────────────────────────────────────────────────
    tb, tf = _txbox(slide, MARGIN, CLAIM_Y, SLIDE_W - 2 * MARGIN, CLAIM_H)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    _run(tf.paragraphs[0], str(data["claim"]), size=28, color=ACCENT)

    # ── left card ────────────────────────────────────────────────────────────
    _rect(slide, MARGIN, BODY_Y, LEFT_W, BODY_H, fill=WHITE, line=INK)
    inner_x = MARGIN + CARD_PAD
    inner_w = LEFT_W - 2 * CARD_PAD

    tb, tf = _txbox(slide, inner_x, BODY_Y + CARD_PAD, inner_w - LOGO_W - 0.12, 0.32)
    _run(tf.paragraphs[0], str(data["name"]), size=15, bold=True)

    tb, tf = _txbox(slide, inner_x, BODY_Y + CARD_PAD + 0.30, inner_w - LOGO_W - 0.12, 0.26)
    metaline = (f"Ort: {data['location']} / Gründung: {data['founded']} / "
                f"Team: {data['team_size']}")
    _run(tf.paragraphs[0], metaline, size=9.5, italic=True)

    logo = data.get("logo")
    logo_x = MARGIN + LEFT_W - CARD_PAD - LOGO_W
    if logo:
        lp = (base_dir / str(logo)).resolve()
        if lp.exists():
            from PIL import Image
            iw, ih = Image.open(lp).size
            scale = min(LOGO_W / (iw / 96), LOGO_H / (ih / 96))
            w_in, h_in = (iw / 96) * scale, (ih / 96) * scale
            slide.shapes.add_picture(
                str(lp), Inches(logo_x + LOGO_W - w_in), Inches(BODY_Y + CARD_PAD),
                Inches(w_in), Inches(h_in))
        else:
            print(f"      ! logo not found: {lp}")

    # Sections share ONE text box on purpose: edited text reflows instead of
    # overflowing a fixed frame, which is the whole point of the .pptx export.
    # Auto-shrink on top of that, so adding a sentence nudges the type down a
    # notch rather than spilling out of the card. The generated pages sit at
    # ~87-91% of the box, so there is room, but not unlimited room.
    tb, tf = _txbox(slide, inner_x, BODY_Y + CARD_PAD + 0.72,
                    inner_w, BODY_H - CARD_PAD * 2 - 0.72)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    sections = data["sections"]
    first = True
    for key, heading in SECTIONS:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        if not first:
            p.space_before = Pt(9)
        _run(p, heading, size=10, bold=True, underline=True)
        pb = tf.add_paragraph()
        pb.space_before = Pt(2)
        _run(pb, " ".join(str(sections[key]).split()), size=9.5)
        first = False

    # ── right column ─────────────────────────────────────────────────────────
    visuals = data["visuals"]
    for (key, default_label), y in zip(VISUALS, (BODY_Y, BODY_Y + VIS_H + VIS_GAP)):
        _visual(slide, visuals.get(key) or {}, default_label, base_dir, tmp_dir, y)


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def main() -> int:
    ap = argparse.ArgumentParser(description="Export GT Hub one-pagers to editable .pptx")
    ap.add_argument("data", nargs="+")
    ap.add_argument("--out-dir", default=None)
    ap.add_argument("--combine", default=None,
                    help="write every page into ONE .pptx at this path instead of one file each")
    args = ap.parse_args()

    tmp_dir = Path(".onepager_tmp")
    combined = new_deck() if args.combine else None
    code = 0

    for raw in args.data:
        path = Path(raw).resolve()
        if not path.exists():
            print(f"  ✗ {raw}: not found")
            code = 1
            continue
        data = yaml.safe_load(path.read_text(encoding="utf-8")) or {}
        problems = validate(data, path)
        if problems:
            print(f"  ✗ {path.name}: {len(problems)} problem(s)")
            for p in problems:
                print(f"      - {p}")
            code = 1
            continue

        if combined is not None:
            build_slide(combined, data, path.parent, tmp_dir)
            print(f"  ✓ {path.name} -> slide added")
        else:
            prs = new_deck()
            build_slide(prs, data, path.parent, tmp_dir)
            out_dir = Path(args.out_dir).resolve() if args.out_dir else path.parent
            out_dir.mkdir(parents=True, exist_ok=True)
            out = out_dir / f"{path.stem}_onepager.pptx"
            prs.save(str(out))
            print(f"  ✓ {path.name} -> {out}")

    if combined is not None:
        out = Path(args.combine).resolve()
        out.parent.mkdir(parents=True, exist_ok=True)
        combined.save(str(out))
        print(f"  ✓ combined deck -> {out}")

    if tmp_dir.exists():
        for f in tmp_dir.iterdir():
            f.unlink()
        tmp_dir.rmdir()
    return code


if __name__ == "__main__":
    sys.exit(main())
