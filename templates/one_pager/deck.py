"""
Pitch-deck parsing for the one-pager generator — PDF and PPTX, text and images.

ISOLATION CONTRACT (see FORMAT.md §7): this module imports NOTHING from the
VC-scouting pipeline. Only stdlib plus fitz (PyMuPDF), python-pptx and Pillow —
all already pinned in requirements.txt. It opens no database, makes no network
call, and reads no project config. tests/test_one_pager_isolation.py enforces
this automatically.

Two jobs:
  parse_deck(path)            -> Deck   (per-slide text, so every drafted fact
                                         can cite a slide number)
  extract_images(deck, out)   -> list   (every embedded image + a full-page
                                         render per slide, for a human to pick
                                         the two the format needs)

Deliberately does NOT choose which images are good. That is visual judgement a
text model does badly: on LIGARO's own material the largest, most prominent
assets were funder and partner logos (EXIST, Leibniz Universität, Future
Greentech Incubator), while the only real product shot was the poster frame of
a site video. Auto-picking would have confidently chosen wrong.
"""
from __future__ import annotations

import io
import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional

logger = logging.getLogger(__name__)

# Downscale/quality for extracted assets. Same treatment export_pptx.py applies
# so a picked image is already sized for both the HTML (base64-embedded) and
# PPTX outputs without a second pass.
MAX_IMAGE_WIDTH = 1500
JPEG_QUALITY = 84

# Below this many characters a "slide" is a section divider or a bare logo wall,
# not content worth feeding a model. Measured against the real ONOX deck, whose
# chapter slides ("Problem", "Lösung", "Markt") carry exactly one word.
MIN_CONTENT_CHARS = 40


@dataclass
class Slide:
    number: int                      # 1-based, matches what a human sees
    text: str = ""
    image_count: int = 0

    @property
    def is_content(self) -> bool:
        return len(self.text.strip()) >= MIN_CONTENT_CHARS


@dataclass
class Deck:
    path: Path
    kind: str                        # "pdf" | "pptx"
    slides: List[Slide] = field(default_factory=list)

    @property
    def full_text(self) -> str:
        """Every slide's text, slide-numbered so a citation stays traceable."""
        return "\n\n".join(
            f"[Folie {s.number}]\n{s.text.strip()}" for s in self.slides if s.text.strip()
        )

    def content_text(self, max_chars: int) -> str:
        """
        Text for the model, capped. Content-dense slides first: a 33-slide deck
        easily exceeds an 8k context, and the chapter dividers carry no
        information, so spending the budget on them would push out real content.
        Slide order is preserved within the selection — the narrative sequence
        (problem -> solution -> market) is itself a signal.
        """
        ranked = sorted(
            (s for s in self.slides if s.is_content),
            key=lambda s: len(s.text), reverse=True,
        )
        chosen, used = [], 0
        for s in ranked:
            block = f"[Folie {s.number}]\n{s.text.strip()}"
            if used + len(block) > max_chars:
                continue
            chosen.append(s.number)
            used += len(block) + 2
        keep = set(chosen)
        return "\n\n".join(
            f"[Folie {s.number}]\n{s.text.strip()}"
            for s in self.slides if s.number in keep
        )


class DeckError(RuntimeError):
    """Unreadable/unsupported deck. The one hard failure — no deck, no job."""


def parse_deck(path) -> Deck:
    p = Path(path).expanduser()
    if not p.exists():
        raise DeckError(f"deck not found: {p}")

    suffix = p.suffix.lower()
    if suffix == ".pdf":
        return _parse_pdf(p)
    if suffix == ".pptx":
        return _parse_pptx(p)
    if suffix == ".ppt":
        # Legacy binary PowerPoint. python-pptx reads only the modern
        # OOXML .pptx; there is no pure-python .ppt reader worth adding as a
        # dependency. Say exactly what to do instead of failing obscurely.
        raise DeckError(
            f"{p.name} is a legacy .ppt file, which cannot be read directly. "
            "Open it in PowerPoint/Keynote and re-save as .pptx, or export to "
            "PDF, then re-run."
        )
    raise DeckError(f"unsupported deck format {suffix!r} — use .pdf or .pptx")


def _parse_pdf(p: Path) -> Deck:
    import fitz  # PyMuPDF

    deck = Deck(path=p, kind="pdf")
    try:
        doc = fitz.open(str(p))
    except Exception as exc:
        raise DeckError(f"could not open {p.name}: {exc}") from exc

    try:
        for i in range(doc.page_count):
            page = doc[i]
            try:
                text = page.get_text() or ""
            except Exception as exc:
                # A single malformed page must not lose the other 32. This is
                # not hypothetical: a real Memminger Zeitung edition had one
                # page whose XObject reference was corrupt in the publisher's
                # own PDF and raised on every text-extraction attempt.
                logger.warning(f"[deck] page {i + 1}: text extraction failed ({exc}) — skipping its text")
                text = ""
            try:
                n_images = len(page.get_images(full=True))
            except Exception:
                n_images = 0
            deck.slides.append(Slide(number=i + 1, text=_clean(text), image_count=n_images))
    finally:
        doc.close()

    if not deck.slides:
        raise DeckError(f"{p.name} has no pages")
    return deck


def _parse_pptx(p: Path) -> Deck:
    from pptx import Presentation

    deck = Deck(path=p, kind="pptx")
    try:
        prs = Presentation(str(p))
    except Exception as exc:
        raise DeckError(f"could not open {p.name}: {exc}") from exc

    for i, slide in enumerate(prs.slides, start=1):
        parts, n_images = [], 0
        for shape in slide.shapes:
            try:
                if shape.shape_type == 13 or getattr(shape, "image", None) is not None:
                    n_images += 1
            except Exception:
                pass
            try:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text.strip())
            except Exception:
                continue
            # Tables carry real content on comparison slides (ONOX's competitor
            # matrix, LIGARO's recyclability table) — losing them would lose
            # exactly the material section 3 needs.
            try:
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
            except Exception:
                pass
        deck.slides.append(Slide(number=i, text=_clean("\n".join(parts)), image_count=n_images))

    if not deck.slides:
        raise DeckError(f"{p.name} has no slides")
    return deck


def _clean(text: str) -> str:
    text = text.replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# ── Image extraction ─────────────────────────────────────────────────────────

def extract_images(deck: Deck, out_dir) -> List[str]:
    """
    Write every embedded image plus one full-page render per slide into
    out_dir, named so a human can map a file straight back to the slide they
    remember: slide07_img2.jpg / slide07_page.jpg.

    Returns the relative filenames written. Never raises — an image-extraction
    failure degrades to "fewer candidates", never to a lost run.
    """
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    written: List[str] = []

    if deck.kind == "pdf":
        written = _extract_pdf_images(deck, out)
    else:
        written = _extract_pptx_images(deck, out)

    if not written:
        logger.warning(f"[deck] no images could be extracted from {deck.path.name}")
    return written


def _save(img_bytes: bytes, dest: Path) -> Optional[str]:
    """Normalise to a downscaled RGB JPEG. Returns the filename, or None."""
    try:
        from PIL import Image

        im = Image.open(io.BytesIO(img_bytes))
        if im.mode in ("RGBA", "LA", "P"):
            im = im.convert("RGBA")
            bg = Image.new("RGB", im.size, (255, 255, 255))
            bg.paste(im, mask=im.split()[-1])
            im = bg
        else:
            im = im.convert("RGB")
        # Skip tiny decorative fragments (bullets, rules, spacer pixels) — they
        # are never a one-pager visual and only add noise to the candidate list.
        if im.width < 120 or im.height < 120:
            return None
        if im.width > MAX_IMAGE_WIDTH:
            im = im.resize(
                (MAX_IMAGE_WIDTH, int(im.height * MAX_IMAGE_WIDTH / im.width)),
                Image.LANCZOS,
            )
        im.save(dest, "JPEG", quality=JPEG_QUALITY, optimize=True)
        return dest.name
    except Exception as exc:
        logger.debug(f"[deck] could not save {dest.name}: {exc}")
        return None


def _extract_pdf_images(deck: Deck, out: Path) -> List[str]:
    import fitz

    written: List[str] = []
    try:
        doc = fitz.open(str(deck.path))
    except Exception as exc:
        logger.warning(f"[deck] reopen for images failed: {exc}")
        return written

    try:
        for i in range(doc.page_count):
            page = doc[i]
            n = i + 1
            try:
                for j, info in enumerate(page.get_images(full=True), start=1):
                    try:
                        pix = fitz.Pixmap(doc, info[0])
                        if pix.n - pix.alpha >= 4:      # CMYK -> RGB
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        name = _save(pix.tobytes("png"), out / f"slide{n:02d}_img{j}.jpg")
                        if name:
                            written.append(name)
                    except Exception as exc:
                        logger.debug(f"[deck] slide {n} image {j}: {exc}")
            except Exception as exc:
                logger.debug(f"[deck] slide {n}: listing images failed: {exc}")

            # A full-page render is the reliable fallback: a slide whose visual
            # is vector art or a composed layout has no single embedded bitmap
            # to pull, but still renders correctly.
            try:
                name = _save(page.get_pixmap(dpi=110).tobytes("png"), out / f"slide{n:02d}_page.jpg")
                if name:
                    written.append(name)
            except Exception as exc:
                logger.debug(f"[deck] slide {n} page render failed: {exc}")
    finally:
        doc.close()
    return written


def _extract_pptx_images(deck: Deck, out: Path) -> List[str]:
    from pptx import Presentation

    written: List[str] = []
    try:
        prs = Presentation(str(deck.path))
    except Exception as exc:
        logger.warning(f"[deck] reopen for images failed: {exc}")
        return written

    for i, slide in enumerate(prs.slides, start=1):
        j = 0
        for shape in slide.shapes:
            try:
                image = getattr(shape, "image", None)
                if image is None:
                    continue
                j += 1
                name = _save(image.blob, out / f"slide{i:02d}_img{j}.jpg")
                if name:
                    written.append(name)
            except Exception as exc:
                logger.debug(f"[deck] slide {i} shape image: {exc}")
    # No page-render fallback for PPTX: rendering a slide needs PowerPoint or
    # LibreOffice, neither of which is a dependency here (and LibreOffice is
    # not installed on this machine). Export the deck to PDF if a slide's
    # visual is a composed layout rather than one embedded picture.
    return written
