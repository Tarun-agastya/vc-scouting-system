"""
One-pager generator — deck parsing, grounding, and YAML assembly.

Pure and hermetic on purpose: no Ollama, no network, no Postgres, no Qdrant.
Real PDF and PPTX files are built in-test with fitz and python-pptx (both
already dependencies) rather than committed as fixtures, so the tests exercise
the actual parsers against actual files while staying self-contained. That
also means these pass with the whole scouting stack switched off — the
condition the isolation exists for.
"""
import sys
from pathlib import Path

import pytest

ONE_PAGER_DIR = Path(__file__).resolve().parent.parent / "templates" / "one_pager"
sys.path.insert(0, str(ONE_PAGER_DIR))

import deck as deck_mod            # noqa: E402
import generate as gen             # noqa: E402
import llm as llm_mod              # noqa: E402


# ── fixtures: real files, built on the fly ───────────────────────────────────

SLIDES = [
    "Problem",                                    # chapter divider — not content
    "Schwankende Dieselkosten belasten Landwirte zunehmend und 65.000 Betriebe "
    "erzeugen bereits eigenen Strom.",
    "Der ONOX 1 ist ein vollelektrischer Traktor mit Wechselmodulen, der ohne "
    "Standzeiten beim Laden den ganzen Tag arbeitet.",
    "Gegruendet 2021 in Isny. Team: 11 Personen. 11t CO2 Einsparung pro Jahr.",
]


@pytest.fixture
def pdf_deck(tmp_path):
    import fitz
    doc = fitz.open()
    for text in SLIDES:
        page = doc.new_page()
        y = 90
        for line in text.split(". "):
            page.insert_text((60, y), line, fontsize=11)
            y += 18
    path = tmp_path / "deck.pdf"
    doc.save(str(path))
    doc.close()
    return path


@pytest.fixture
def pptx_deck(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    for text in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(3))
        box.text_frame.text = text
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


# ── parsing: both formats reach the same shape ───────────────────────────────

def test_pdf_and_pptx_parse_to_the_same_shape(pdf_deck, pptx_deck):
    a, b = deck_mod.parse_deck(pdf_deck), deck_mod.parse_deck(pptx_deck)
    assert (a.kind, b.kind) == ("pdf", "pptx")
    assert len(a.slides) == len(b.slides) == len(SLIDES)
    for d in (a, b):
        assert "ONOX 1" in d.full_text
        assert "[Folie 3]" in d.full_text          # slide numbers survive, for citations


def test_chapter_dividers_are_not_treated_as_content(pdf_deck):
    d = deck_mod.parse_deck(pdf_deck)
    assert d.slides[0].text.strip() == "Problem"
    assert d.slides[0].is_content is False, "a one-word chapter slide is not content"
    assert d.slides[1].is_content is True


def test_content_text_respects_its_budget(pdf_deck):
    d = deck_mod.parse_deck(pdf_deck)
    assert len(d.content_text(120)) <= 200          # budget honoured (+ slide headers)
    assert d.content_text(5000)                     # a generous budget still returns text


def test_unreadable_and_unsupported_decks_fail_clearly(tmp_path):
    with pytest.raises(deck_mod.DeckError, match="not found"):
        deck_mod.parse_deck(tmp_path / "nope.pdf")

    legacy = tmp_path / "old.ppt"
    legacy.write_bytes(b"\xd0\xcf\x11\xe0")
    with pytest.raises(deck_mod.DeckError, match="re-save as .pptx|legacy"):
        deck_mod.parse_deck(legacy)

    other = tmp_path / "x.key"
    other.write_bytes(b"x")
    with pytest.raises(deck_mod.DeckError, match="unsupported"):
        deck_mod.parse_deck(other)


def test_image_extraction_never_raises_and_names_by_slide(pdf_deck, tmp_path):
    d = deck_mod.parse_deck(pdf_deck)
    names = deck_mod.extract_images(d, tmp_path / "assets")
    assert names, "a PDF should always yield at least page renders"
    assert all(n.startswith("slide") for n in names)
    assert any("_page" in n for n in names)


# ── grounding ────────────────────────────────────────────────────────────────

SOURCE = ("11t CO2 pro Jahr. Gesamtkosten nach 8 Jahren: ONOX 95.861 EUR gegenueber "
          "Verbrenner 121.300 EUR. 600 Newsletter Abos, 25 Probefahrten. Team: 11 Personen.")


def test_grounding_catches_a_number_the_model_computed_itself():
    """
    The real case, from the first live run against an ONOX-shaped deck: given
    121.300 and 95.861, the model wrote "25.439 EUR gesenkt" — arithmetically
    correct, but that figure appears nowhere in the source. A reader would
    assume it came from the deck.
    """
    assert gen._unsupported_numbers("Nach 8 Jahren 25.439 EUR gesenkt.", SOURCE) == ["25.439"]


def test_grounding_passes_numbers_that_are_really_in_the_deck():
    assert gen._unsupported_numbers("11t CO2, 95.861 EUR nach 8 Jahren.", SOURCE) == []


def test_grounding_ignores_thousand_separator_style():
    """95861 and 95.861 are the same number — a format difference must not
    read as a fabrication."""
    assert gen._unsupported_numbers("Kosten 95861 EUR.", SOURCE) == []


def test_grounding_never_touches_prose_without_numbers():
    assert gen._unsupported_numbers("Ein vollelektrischer Traktor für Höfe.", SOURCE) == []


def test_unsupported_team_size_becomes_none_not_a_guess():
    assert gen._supported_meta("11", SOURCE) == "11"
    assert gen._supported_meta("6", SOURCE) is None     # -> rendered as "k. A."
    assert gen._supported_meta(None, SOURCE) is None


# ── YAML assembly ────────────────────────────────────────────────────────────

def _build(pdf_deck, drafted, images=("slide01_page.jpg",), llm_note=None):
    return gen.build_yaml(
        name="ONOX", slug="onox", drafted=drafted,
        deck_obj=deck_mod.parse_deck(pdf_deck), images=list(images),
        url=None, url_ok=False, llm_note=llm_note,
    )


def test_yaml_has_every_field_the_renderer_requires(pdf_deck):
    data = _build(pdf_deck, {
        "claim": "Elektrischer Traktor mit Wechselbatterien",
        "location": "Isny", "founded": "2021", "team_size": "11",
        "loesung": "Ein Traktor.", "mehrwerte": "11t CO2 pro Jahr.",
        "usp": "Wechselmodule.", "zielgruppe": "Höfe.", "geschaeftsmodell": "Verkauf.",
    })
    for key in ("claim", "name", "location", "founded", "team_size", "sections", "visuals"):
        assert key in data
    assert set(data["sections"]) == set(gen.SECTION_LABELS)
    assert data["review"]["status"] == "draft", "generated pages are never pre-approved"


def test_a_failed_draft_still_produces_a_complete_yaml(pdf_deck):
    """Ollama being down must cost the prose, never the run."""
    data = _build(pdf_deck, None, llm_note="Ollama nicht erreichbar")
    assert data["claim"] == ""
    assert all(v == "" for v in data["sections"].values())
    assert data["location"] == "k. A."
    q = " ".join(data["review"]["open_questions"])
    assert "Ollama" in q
    for heading in gen.SECTION_LABELS.values():
        assert heading in q, f"the human must be told '{heading}' needs writing"


def test_missing_images_are_flagged_rather_than_faked(pdf_deck):
    data = _build(pdf_deck, None, images=())
    assert "placeholder" in data["visuals"]["visual_solution"]
    assert "image" not in data["visuals"]["visual_solution"]
    assert any("keine Bilder" in q.lower() or "manuell" in q.lower()
               for q in data["review"]["open_questions"])


def test_sources_cite_the_deck_and_the_slides_used(pdf_deck):
    data = _build(pdf_deck, None)
    joined = " ".join(data["sources"])
    assert "deck.pdf" in joined and "Folien" in joined
    assert "Inhaltsfolien ausgewertet" in joined


# ── llm helpers (no network) ─────────────────────────────────────────────────

def test_llm_normalise_strips_sentinels_claim_period_and_team_noise():
    out = llm_mod._normalise({
        "claim": "Elektrischer Traktor mit Wechselbatterien.",   # FORMAT.md forbids the period
        "location": "Isny", "founded": "2021", "team_size": "11 Personen",
        "loesung": "Ein Traktor.", "mehrwerte": "", "usp": "   ",
        "zielgruppe": "Höfe.", "geschaeftsmodell": "Verkauf.",
    })
    assert out["claim"] == "Elektrischer Traktor mit Wechselbatterien"
    assert out["team_size"] == "11"
    assert out["mehrwerte"] is None and out["usp"] is None   # "" means "deck doesn't say"


def test_llm_draft_returns_none_when_ollama_is_unreachable(monkeypatch):
    monkeypatch.setattr(llm_mod, "BASE_URL", "http://localhost:1")
    monkeypatch.setattr(llm_mod.time, "sleep", lambda *_: None)   # don't pay the retry pause
    assert llm_mod.draft("X", "[Folie 1]\nEin Traktor.") is None


def test_slugify_handles_umlauts_and_punctuation():
    assert gen.slugify("Hula Earth") == "hula_earth"
    assert gen.slugify("Müller & Söhne GmbH") == "mueller_soehne_gmbh"
    assert gen.slugify("!!!") == "startup"
